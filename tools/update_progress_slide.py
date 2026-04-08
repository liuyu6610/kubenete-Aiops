# -*- coding: utf-8 -*-
"""给毕业设计 PPT 幂等追加“当前进展与开发记录”页。"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "docs" / "KubeSentinel-thesis-defense.pptx"
SLIDE_TITLE = "当前进展与开发记录"

NAVY = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(8, 145, 178)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(234, 88, 12)
RED = RGBColor(220, 38, 38)
SLATE = RGBColor(71, 85, 105)
WHITE = RGBColor(255, 255, 255)


def set_text(frame, text, size=16, bold=False, color=NAVY, align=None):
    """统一文本样式，保持与原 PPT 风格接近。"""
    frame.clear()
    paragraph = frame.paragraphs[0]
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.45))
    set_text(title_box.text_frame, title, 22, True, NAVY)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.78), Inches(1.35), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(9.2), Inches(0.25), Inches(3.5), Inches(0.35))
        set_text(sub.text_frame, subtitle, 9, False, SLATE, PP_ALIGN.RIGHT)


def add_footer(slide, page):
    footer = slide.shapes.add_textbox(Inches(0.55), Inches(7.13), Inches(11.9), Inches(0.2))
    set_text(footer.text_frame, f"KubeSentinel LLM AIOps 毕业设计 | {page:02d}", 8, False, SLATE, PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, title, body, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(203, 213, 225)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.12), Inches(w - 0.32), Inches(0.3))
    set_text(title_box.text_frame, title, 13, True, NAVY)
    body_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.5), Inches(w - 0.32), Inches(h - 0.62))
    set_text(body_box.text_frame, body, 10, False, SLATE)


def slide_text(slide):
    """提取 slide 文本，用于识别旧的进展页。"""
    chunks = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            chunks.append(shape.text)
    return "\n".join(chunks)


def remove_progress_slides(prs):
    """删除已存在的同名进展页，使脚本可重复执行。"""
    slide_ids = prs.slides._sldIdLst  # python-pptx 暂无公开删除 API
    for index in range(len(prs.slides) - 1, -1, -1):
        if SLIDE_TITLE in slide_text(prs.slides[index]):
            rel_id = slide_ids[index].rId
            prs.part.drop_rel(rel_id)
            del slide_ids[index]


def main():
    if not PPTX.exists():
        raise FileNotFoundError(f"找不到 PPT 文件: {PPTX}")

    prs = Presentation(PPTX)
    remove_progress_slides(prs)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_title(slide, SLIDE_TITLE, "Progress")

    add_card(
        slide,
        0.75,
        1.15,
        3.75,
        1.3,
        "已完成：仓库阅读与架构梳理",
        "已通读前端、后端、Operator、监控告警、n8n 审批与脚本目录，明确系统从告警接入到自愈执行的闭环。",
        BLUE,
    )
    add_card(
        slide,
        4.8,
        1.15,
        3.75,
        1.3,
        "已完成：答辩材料准备",
        "已生成架构图与答辩 PPT，并将项目定位、核心模块、闭环流程和风险点整理为可展示内容。",
        GREEN,
    )
    add_card(
        slide,
        8.85,
        1.15,
        3.75,
        1.3,
        "已完成：基础工程校验",
        "已验证 Python 后端编译、Go Operator 测试与前端构建路径，为后续服务器部署打底。",
        CYAN,
    )

    add_card(
        slide,
        0.75,
        3.05,
        3.75,
        1.45,
        "进行中：租服务器",
        "准备租用云服务器搭建 Kubernetes 集群，后续部署 Prometheus、Alertmanager、Loki、后端、前端、n8n 与 Operator。",
        ORANGE,
    )
    add_card(
        slide,
        4.8,
        3.05,
        3.75,
        1.45,
        "进行中：真实环境运行",
        "计划在服务器集群中复现实验链路：制造告警、AI 分析、审批回调、创建 HealingRule、Operator 执行修复。",
        RED,
    )
    add_card(
        slide,
        8.85,
        3.05,
        3.75,
        1.45,
        "待交付：开发过程材料",
        "后续整理环境搭建、部署步骤、问题记录、测试结果与演示截图，作为毕业设计开发过程佐证。",
        NAVY,
    )

    note = slide.shapes.add_textbox(Inches(1.0), Inches(5.25), Inches(11.25), Inches(0.9))
    set_text(
        note.text_frame,
        "这页只作为答辩中的开发记录摘要，不单独生成开发文档；正式文档会在服务器部署与联调完成后补齐。",
        15,
        True,
        NAVY,
        PP_ALIGN.CENTER,
    )
    add_footer(slide, len(prs.slides))
    prs.save(PPTX)
    print(f"{PPTX} slides={len(prs.slides)}")


if __name__ == "__main__":
    main()
